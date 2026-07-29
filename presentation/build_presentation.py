#!/usr/bin/env python3
"""Generate the 5-minute jury deck (pptx) in the Lamer Konekte visual identity."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0B, 0x25, 0x45)
TURQ = RGBColor(0x1B, 0x9A, 0xAA)
CORAL = RGBColor(0xFF, 0x6B, 0x57)
FOAM = RGBColor(0xDE, 0xF2, 0xF1)
WHITE = RGBColor(0xFA, 0xF9, 0xF6)

HERE = Path(__file__).resolve().parent
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(title: str, bullets: list[str], accent=TURQ, title_color=WHITE, dark=True, kicker: str | None = None):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY if dark else WHITE
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    if kicker:
        kb = s.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12), Inches(0.5)).text_frame
        kb.text = kicker
        kb.paragraphs[0].runs[0].font.size = Pt(14)
        kb.paragraphs[0].runs[0].font.color.rgb = accent
        kb.paragraphs[0].runs[0].font.bold = True
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.85), Inches(12.1), Inches(1.2)).text_frame
    tb.text = title
    r = tb.paragraphs[0].runs[0]
    r.font.size = Pt(40); r.font.bold = True
    r.font.color.rgb = title_color if dark else NAVY
    body = s.shapes.add_textbox(Inches(0.7), Inches(2.15), Inches(12), Inches(4.9)).text_frame
    body.word_wrap = True
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.space_after = Pt(14)
        f = p.runs[0].font
        f.size = Pt(20)
        f.color.rgb = FOAM if dark else NAVY
        if b.startswith("!"):
            p.runs[0].text = b[1:]
            f.color.rgb = CORAL
            f.bold = True
    return s


# 1 title
s = prs.slides.add_slide(BLANK)
s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
wave = s.shapes.add_shape(1, 0, Inches(5.4), prs.slide_width, Inches(2.1))
wave.fill.solid(); wave.fill.fore_color.rgb = TURQ; wave.line.fill.background()
t = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.6)).text_frame
t.text = "Lamer Konekte"
t.paragraphs[0].runs[0].font.size = Pt(66); t.paragraphs[0].runs[0].font.bold = True
t.paragraphs[0].runs[0].font.color.rgb = WHITE
p = t.add_paragraph(); p.text = "Lapes pli konekte. Desizion pli informe."
p.runs[0].font.size = Pt(28); p.runs[0].font.color.rgb = CORAL; p.runs[0].font.italic = True
p2 = t.add_paragraph(); p2.text = "Team Ctrl200 · Gemma 4 Hackathon · Multimodal Track · Blue Economy"
p2.runs[0].font.size = Pt(18); p2.runs[0].font.color.rgb = FOAM

# 2-10
slide("A fisher's morning, unconnected", [
    "Paper catch declarations → national statistics arrive late and incomplete",
    "Marine forecasts scattered across sources not written for fishers",
    "Regulations like the octopus closure reach the lagoon second-hand",
    "!And almost none of it exists in Morisyen — the language fishers speak",
], kicker="0:00 — THE PROBLEM")

slide("The Blue Economy's invisible actors", [
    "Artisanal fishing feeds Mauritius, but its smallest actors are invisible to its data",
    "No structured catch records → no evidence base for sustainable management",
    "Digitisation keeps failing when it isn't built in the fisher's language",
], kicker="0:25 — NATIONAL BOTTLENECK")

slide("Why Gemma 4", [
    "One model — gemma-4-26b-a4b-it via the official google-genai SDK",
    "Photo understanding + Morisyen + structured JSON + native function calling",
    "Candidate shortlist retrieval: Gemma suggests ONLY from allowed species",
    "Gemma family scales down (E2B/E4B) → a credible on-device roadmap for offline lagoons",
    "!Bounded by design: suggest, never declare · no legality · no invented rules",
], kicker="0:50 — WHY GEMMA")

slide("Moment 1 — Before the trip", [
    "“Ki kalite lamer ena dan Grand Baie zordi?”",
    "Gemma selects get_marine_conditions → Open-Meteo waves, swell, sea temperature",
    "Function trace with latency shown on the Technical Proof page",
    "!Mandatory disclaimer: informational only — confirm official advisories",
], kicker="1:15 — LIVE DEMO · FUNCTION CALLING")

slide("Moment 2 — On the water", [
    "Photo → quality gate (blurry photos never spend tokens)",
    "Constrained suggestion + visible characteristics + honest confidence",
    "Fisher CONFIRMS or corrects — mandatory, always",
    "Measured length with a ruler → deterministic, source-attributed rule check",
    "!29 July: no closure · simulated 1 September (badged): closed season, 2016 regulations, ‘provisional’",
], kicker="2:00 — LIVE DEMO · CATCH FLOW")

slide("Moment 3 — Back ashore", [
    "Catch log + today report",
    "Offline-first: IndexedDB queue syncs when connectivity returns",
    "Declaration draft → PDF → demonstration receipt",
    "!Clearly labelled MOCK ministry endpoint — never presented as real government",
], kicker="3:10 — LIVE DEMO · LOG & DECLARATION")

slide("Engineering proof", [
    "12 allow-listed functions · Pydantic-validated · explicit dispatch · redacted traces",
    "44 backend tests: rule boundaries, prompt injection, privacy, secret hygiene",
    "Prototype benchmark: 32 Morisyen cases · 100% schema validity · 0 safety failures",
    "Training: QLoRA notebooks + leakage-safe data, push-button on Kaggle — status reported honestly",
    "!Provider badge on every response: hosted / mock / local, with real_inference flag",
], kicker="3:50 — TECHNICAL PROOF")

slide("Morisyen-first, honestly", [
    "Entire UI in Kreol Morisien and English (~90 strings each)",
    "Every analysis answers in both languages",
    "Species names marked provisional until native-speaker review — never guessed silently",
    "Licensed dataset: 60 iNaturalist photos, 5 Mauritian species, full attribution",
], kicker="4:20 — MORISYEN")

slide("From paper to policy-grade data", [
    "A working preview of digital catch declarations for the Blue Economy",
    "Humans confirm everything that matters; deterministic code owns the law",
    "Next: native-speaker review · ministry schema consultation · Gemma E2B edge pilot",
    "!Limitation on every screen: verify suggestions and rules with official sources",
], accent=CORAL, kicker="4:45 — IMPACT")

prs.save(str(HERE / "Lamer_Konekte_5_Minute.pptx"))
print("saved", HERE / "Lamer_Konekte_5_Minute.pptx", f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
